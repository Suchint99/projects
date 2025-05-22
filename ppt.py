from pptx import Presentation
import openai

# Function to extract text from PPT slides
def extract_text_from_ppt(ppt_file):
    prs = Presentation(ppt_file)
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to generate simple summaries using OpenAI (replace this with your API key)
def generate_simple_summary(text):
    openai.api_key = "YOUR_OPENAI_API_KEY"  # Add your OpenAI API key here
    response = openai.Completion.create(
        engine="text-davinci-003",
        prompt=f"Summarize the following text in simple words:\n\n{text}",
        max_tokens=100,  # Adjust this based on desired summary length
        temperature=0.5
    )
    return response.choices[0].text.strip()

# Function to summarize each slide and save results
def summarize_ppt(ppt_file, output_file="ppt_summary.txt"):
    # Extract text from PowerPoint
    text = extract_text_from_ppt(ppt_file)
    print("Text extracted from PowerPoint.")

    # Generate summary using AI
    summary = generate_simple_summary(text)
    print("Summary generated.")

    # Save summary to a text file
    with open(output_file, "w") as file:
        file.write("Summary of PowerPoint:\n")
        file.write(summary)
    print(f"Summary saved to {output_file}")

# Example usage
ppt_file = "First_neural_Network-Model.pptx"  # Replace with the path to your .pptx file
summarize_ppt(ppt_file)
